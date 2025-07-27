"""
Office特定格式解析器
提供Excel和PowerPoint的专门解析功能
"""

from typing import Optional

from ...utils import get_logger
from .image_extractor import OfficeImageExtractor, PowerPointTextAnalyzer


class ExcelParser:
    """Excel专用解析器"""

    def __init__(self):
        """初始化Excel解析器"""
        self.logger = get_logger("excel_parser")
        self.image_extractor = OfficeImageExtractor()

    def parse_to_markdown(self, file_path: str) -> str:
        """
        使用 pandas 将 Excel 文件转换为 Markdown 表格格式
        
        Args:
            file_path: Excel文件路径
            
        Returns:
            Markdown格式的内容
        """
        try:
            self.logger.info(f"使用pandas转换Excel文件为Markdown: {file_path}")
            
            try:
                import pandas as pd
            except ImportError:
                raise Exception("缺少pandas库，请安装: pip install pandas")
            
            # 读取Excel文件的所有工作表
            excel_file = pd.ExcelFile(file_path)
            markdown_parts = []
            
            for sheet_name in excel_file.sheet_names:
                self.logger.debug(f"处理工作表: {sheet_name}")
                
                # 读取工作表数据
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # 跳过空工作表
                if df.empty:
                    continue
                
                # 添加工作表标题
                markdown_parts.append(f"## 工作表: {sheet_name}")
                markdown_parts.append("")
                
                # 转换为Markdown表格
                markdown_table = df.to_markdown(index=False, tablefmt='github')
                markdown_parts.append(markdown_table)
                markdown_parts.append("")
            
            if markdown_parts:
                markdown_content = "\n".join(markdown_parts)
                self.logger.info(f"Excel转Markdown成功，处理了 {len(excel_file.sheet_names)} 个工作表")
                return markdown_content
            else:
                self.logger.warning("Excel文件中没有有效数据")
                return "# Excel文件\n\n*文件中没有有效数据*"
                
        except Exception as e:
            self.logger.error(f"Excel转Markdown失败: {e}")
            raise Exception(f"Excel解析失败: {e}")

    def extract_images(self, file_path: str, temp_image_dir: str) -> None:
        """
        提取Excel文件中的图片
        
        Args:
            file_path: Excel文件路径
            temp_image_dir: 临时图片保存目录
        """
        self.image_extractor.extract_excel_images(file_path, temp_image_dir)


class PowerPointParser:
    """PowerPoint专用解析器"""

    def __init__(self):
        """初始化PowerPoint解析器"""
        self.logger = get_logger("powerpoint_parser")
        self.image_extractor = OfficeImageExtractor()
        self.text_analyzer = PowerPointTextAnalyzer()

    def parse_to_markdown(self, file_path: str, temp_image_dir: str = None) -> str:
        """
        使用 python-pptx 将 PowerPoint 文件转换为 Markdown 格式，并提取图片
        
        Args:
            file_path: PowerPoint文件路径
            temp_image_dir: 临时图片保存目录
            
        Returns:
            Markdown格式的内容
        """
        try:
            self.logger.info(f"使用python-pptx转换PowerPoint为Markdown: {file_path}")
            
            try:
                from pptx import Presentation
            except ImportError:
                raise Exception("缺少python-pptx库，请安装: pip install python-pptx")
            
            # 加载演示文稿
            prs = Presentation(file_path)
            markdown_parts = []
            
            # 添加文档标题
            markdown_parts.append("# PowerPoint 演示文稿")
            markdown_parts.append("")
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                self.logger.debug(f"处理幻灯片: {slide_idx}")
                
                # 添加幻灯片标题
                markdown_parts.append(f"## 幻灯片 {slide_idx}")
                markdown_parts.append("")
                
                # 提取幻灯片中的所有文本和图片
                slide_texts = []
                slide_images = []
                
                for shape in slide.shapes:
                    # 处理文本内容
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        # 简单的格式判断和转换
                        if self.text_analyzer.is_title_text(shape, text):
                            # 标题文本
                            slide_texts.append(f"### {text}")
                        elif self.text_analyzer.is_bullet_text(text):
                            # 列表文本
                            lines = text.split('\n')
                            for line in lines:
                                if line.strip():
                                    slide_texts.append(f"- {line.strip()}")
                        else:
                            # 普通文本
                            slide_texts.append(text)
                
                # 提取图片
                if temp_image_dir:
                    slide_images = self.image_extractor.extract_pptx_images_from_slide(
                        slide, slide_idx - 1, temp_image_dir
                    )
                
                # 添加文本内容
                if slide_texts:
                    markdown_parts.extend(slide_texts)
                else:
                    markdown_parts.append("*此幻灯片无文本内容*")
                
                # 添加图片引用
                if slide_images:
                    markdown_parts.append("")
                    for img_ref in slide_images:
                        markdown_parts.append(img_ref)
                
                markdown_parts.append("")  # 幻灯片间空行
            
            if len(markdown_parts) > 2:  # 除了标题还有内容
                markdown_content = "\n".join(markdown_parts)
                self.logger.info(f"PowerPoint转Markdown成功，处理了 {len(prs.slides)} 个幻灯片")
                return markdown_content
            else:
                self.logger.warning("PowerPoint文件中没有有效内容")
                return "# PowerPoint 演示文稿\n\n*演示文稿中没有有效内容*"
                
        except Exception as e:
            self.logger.error(f"PowerPoint转Markdown失败: {e}")
            raise Exception(f"PowerPoint解析失败: {e}") 